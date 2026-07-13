#!/usr/bin/env python3
"""
Build the "Thinking in P2: Functional Decomposition" talk deck as a .pptx
that imports cleanly into Keynote (File > Open), with presenter notes on
every slide.

Source of truth: manuals/p2-architect-guide/opus-master/architect-guide-body.md
(Part I intro + Part II, Chapters 5-9). Content and language are drawn from
that manual; nothing here is invented.

Regenerate:  python3 build_deck.py
Output:      thinking-in-p2-decomposition.pptx  (next to this script)

Keynote is theme-neutral here on purpose: import the .pptx, then apply a
Keynote theme. Slides carry concise bullets; the talk track lives in the
presenter notes.
"""

import os
import re
import struct
import zipfile

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = __file__.rsplit("/", 1)[0]

# ---- Brand (Iron Sheep Productions) -----------------------------------------
# Maroon sampled from the logo wordmark/traces. Drives the accent colour so the
# deck reads as Iron Sheep's, not generic. Background is warm-neutral (not the
# earlier cool blue, which fought the maroon logo).
MAROON = RGBColor(0x6C, 0x0D, 0x0E)
MAROON_HEX = "6C0D0E"
LOGO_FULL = HERE + "/company-logo/logo-full-transparent.png"   # wordmark, white-keyed transparent
LOGO_MARK = HERE + "/company-logo/sheep-mark-transparent.png"  # sheep bug, white-keyed transparent

# ---- Footers / title-slide credit (standard talk conventions) ---------------
# Title slide carries name + company + date; content slides carry a muted
# company footer (bottom-left) and a live slide-number field (bottom-right).
# No date in the running footers. Muted gray so both recede under the content.
PRESENTER = "Stephen M. Moraco"
TALK_DATE = "July 8, 2026"
CREDIT_COLOR = RGBColor(0x44, 0x44, 0x44)  # readable dark gray
FOOT_PT = 12


def _slidenum_field(text_frame, number):
    """Insert a live PowerPoint slide-number field, coloured brand maroon. Its
    literal fallback text is set to the real position, so even an importer that
    won't evaluate the field still shows the correct number."""
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    pel = p._p
    fld = pel.makeelement(qn("a:fld"), {
        "id": "{5F3D8B12-1A2B-4C3D-9E4F-%012d}" % number, "type": "slidenum"})
    rPr = pel.makeelement(qn("a:rPr"), {"lang": "en-US", "sz": str(FOOT_PT * 100)})
    fill = pel.makeelement(qn("a:solidFill"), {})
    clr = pel.makeelement(qn("a:srgbClr"), {"val": MAROON_HEX})
    fill.append(clr)
    rPr.append(fill)
    fld.append(rPr)
    t = pel.makeelement(qn("a:t"), {})
    t.text = str(number)
    fld.append(t)
    pel.append(fld)


def add_footer(slide, number):
    """Sheep brand-mark (bottom-left) + live maroon slide number (bottom-right)."""
    slide.shapes.add_picture(LOGO_MARK, Inches(0.6), Inches(6.82), height=Inches(0.52))
    num = slide.shapes.add_textbox(Inches(11.13), Inches(7.02), Inches(1.6), Inches(0.4))
    _slidenum_field(num.text_frame, number)


def color_title(slide):
    """Set a content-slide title in brand maroon (the deck's title accent)."""
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = MAROON

# Subtle diagonal background gradient for a bit of color. Kept LIGHT so the
# deck's dark title/body text stays readable. Two hex stops + an angle — tweak
# freely. Applied per-slide (most reliable for Keynote import). If a Keynote
# theme is applied after import, its background may override these.
BG_TOP = RGBColor(0xF5, 0xF1, 0xEF)     # warm off-white (top-left) — harmonizes with maroon
BG_BOTTOM = RGBColor(0xFF, 0xFF, 0xFF)  # white (bottom-right)
BG_ANGLE = 45.0


def apply_gradient_bg(slide):
    fill = slide.background.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = BG_TOP
    stops[1].position = 1.0
    stops[1].color.rgb = BG_BOTTOM
    try:
        fill.gradient_angle = BG_ANGLE
    except (NotImplementedError, ValueError):
        pass


def wire_notes_master(path):
    """Keynote fix. python-pptx creates the notesMaster part and its
    relationship but omits the <p:notesMasterIdLst> element from
    presentation.xml. PowerPoint and LibreOffice follow the relationship
    regardless; Keynote's strict importer rejects any notes-bearing package
    that lacks the element ('invalid file format'). Post-process the saved
    file to insert it in the schema-correct position (right after
    </p:sldMasterIdLst>), using the notesMaster relationship id. Verified: with
    this element Keynote opens the deck; without it, it will not.
    """
    zin = zipfile.ZipFile(path, "r")
    rels = zin.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
    m = re.search(r'Id="([^"]+)"[^>]*relationships/notesMaster"', rels)
    pres = zin.read("ppt/presentation.xml").decode("utf-8")
    if m and "notesMasterIdLst" not in pres:
        inject = ('<p:notesMasterIdLst><p:notesMasterId r:id="%s"/>'
                  '</p:notesMasterIdLst>' % m.group(1))
        pres = pres.replace("</p:sldMasterIdLst>",
                            "</p:sldMasterIdLst>" + inject, 1)
    items = [(i, zin.read(i.filename)) for i in zin.infolist()]
    zin.close()
    zout = zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED)
    for item, data in items:
        if item.filename == "ppt/presentation.xml":
            data = pres.encode("utf-8")
        zout.writestr(item, data)
    zout.close()

# Diagnostic: set DECK_NO_IMAGES=1 to build an image-free variant (figure
# slides become plain text placeholders) written to a *-noimg-test.pptx file.
# Used to isolate whether embedded PNGs are what a strict importer (Keynote)
# rejects. Normal runs are unaffected.
NO_IMAGES = os.environ.get("DECK_NO_IMAGES") == "1"


def png_size(path):
    with open(path, "rb") as f:
        head = f.read(24)
    return struct.unpack(">II", head[16:24])  # (width, height)


prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)

# Keynote is strict: overriding the size leaves python-pptx's default
# type="screen4x3" on <p:sldSz>, contradicting the 16:9 cx/cy. PowerPoint and
# LibreOffice ignore the stale attribute; Keynote rejects the file as invalid.
# Drop it so the slide size reads as a plain custom (16:9) size.
_sldSz = prs._element.find(qn("p:sldSz"))
if _sldSz is not None and "type" in _sldSz.attrib:
    del _sldSz.attrib["type"]

TITLE_ONLY = prs.slide_layouts[5]   # title only (we add our own body box)
BULLET = prs.slide_layouts[1]       # title + content

# The built-in template's placeholders are sized for the old 4:3 slide
# (0.5in left, 9.0in wide), which leaves a ~3.8in right margin on our 16:9
# slide. Re-place title/body across the full width with symmetric ~0.6in
# margins so content fills the slide evenly.
CONTENT_LEFT = Inches(0.6)
CONTENT_WIDTH = Inches(12.13)        # 13.333 - 2 * 0.6


def _place(shape, left, top, width, height):
    shape.left, shape.top, shape.width, shape.height = left, top, width, height


def bullet_slide(title, bullets, notes):
    """bullets: list of (text, level) or plain strings (level 0)."""
    slide = prs.slides.add_slide(BULLET)
    slide.shapes.title.text = title
    _place(slide.shapes.title, CONTENT_LEFT, Inches(0.3), CONTENT_WIDTH, Inches(1.25))
    color_title(slide)
    body_ph = slide.placeholders[1]
    _place(body_ph, CONTENT_LEFT, Inches(1.6), CONTENT_WIDTH, Inches(5.4))
    body = body_ph.text_frame
    body.word_wrap = True
    first = True
    for b in bullets:
        text, level = (b if isinstance(b, tuple) else (b, 0))
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.text = text
        p.level = level
        for run in p.runs:
            run.font.size = Pt(22 - 2 * min(level, 2))
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def image_slide(title, image_file, notes):
    """A dedicated figure slide: title on top, image centered below.
    Scales the PNG to fit the area under the title, preserving aspect."""
    slide = prs.slides.add_slide(TITLE_ONLY)
    slide.shapes.title.text = title
    _place(slide.shapes.title, CONTENT_LEFT, Inches(0.3), CONTENT_WIDTH, Inches(1.25))
    color_title(slide)
    if NO_IMAGES:
        box = slide.shapes.add_textbox(Inches(1.0), Inches(3.2),
                                       Inches(11.3), Inches(1.0))
        p = box.text_frame.paragraphs[0]
        p.text = "[ figure omitted — image-free load test ]"
        p.runs[0].font.size = Pt(24)
        slide.notes_slide.notes_text_frame.text = notes
        return slide
    w_px, h_px = png_size(HERE + "/" + image_file)
    aspect = w_px / h_px
    # Title placeholder bottom sits at 1.55"; start the image area below it so a
    # tall figure (e.g. the robot map) can't ride up under the title. Height is
    # capped to stay clear of the footer band (~6.95").
    area_w, area_h = 12.4, 5.3           # inches available below the title
    area_top = 1.75
    disp_w = min(area_w, area_h * aspect)
    disp_h = disp_w / aspect
    left = (13.333 - disp_w) / 2
    top = area_top + (area_h - disp_h) / 2
    slide.shapes.add_picture(HERE + "/" + image_file,
                             Inches(left), Inches(top),
                             width=Inches(disp_w), height=Inches(disp_h))
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def title_slide(title, subtitle, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Full brand logo (wordmark included) centered near the top.
    logo_h = 1.7
    logo_w = logo_h * (200 / 269)                 # preserve aspect
    slide.shapes.add_picture(LOGO_FULL, Inches((13.333 - logo_w) / 2), Inches(0.45),
                             height=Inches(logo_h))
    slide.shapes.title.text = title
    _place(slide.shapes.title, CONTENT_LEFT, Inches(2.55), CONTENT_WIDTH, Inches(1.5))
    sub = slide.placeholders[1]
    sub.text = subtitle
    _place(sub, CONTENT_LEFT, Inches(4.05), CONTENT_WIDTH, Inches(0.9))
    # Credit: presenter + date (the company name lives in the logo above).
    credit = slide.shapes.add_textbox(CONTENT_LEFT, Inches(5.2), CONTENT_WIDTH, Inches(1.0))
    first = True
    for text, sz in ((PRESENTER, 18), (TALK_DATE, 14)):
        p = credit.text_frame.paragraphs[0] if first else credit.text_frame.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.color.rgb = CREDIT_COLOR
    slide.notes_slide.notes_text_frame.text = notes
    return slide


# ---------------------------------------------------------------- Slide 1
title_slide(
    "Thinking in P2: Functional Decomposition",
    "The P2 Architect's Guide  •  Part II  •  just released v1.0.0",
    "OPENING (≈60s). Welcome. Last night we shipped the P2 Architect's Guide, "
    "v1.0.0. Tonight I want to give you the heart of it — not the whole book, but "
    "its center: how you decide what goes on which cog in the first place. The book has "
    "three parts — getting a project off the ground, thinking in P2 (decomposition), "
    "and doing that same work with an AI agent. I'll set up Part I in one slide, spend "
    "our time in Part II, then open the floor — because a lot of you have decomposed "
    "real P2 designs, and I want those experiences to feed back into this document. "
    "Plan: ~12 minutes, then questions, then let's compare notes.")

# ---------------------------------------------------------------- Slide 2
bullet_slide(
    "First: Getting a Project Off the Ground (Part I)",
    [
        "Four phases before any cog is assigned:",
        ("Decide what to build — feasibility, parts, pins, what the P2 should NOT do", 1),
        ("Learn the hardware — datasheets, level shifters, the logic analyzer", 1),
        ("Build the capability — interface design, translation, performance, characterize", 1),
        ("Finish & ship — document, post, announce", 1),
        "Output handed to Part II: a wired-up app with a pin map, parts that talk,",
        ("and a feel for their rates and deadlines — the material we carve around", 1),
        "Not one cog is assigned yet. That is the point.",
    ],
    "PART I SETUP (≈90s). Before you can ask Part II's question — which cog owns "
    "what? — there's a whole front end every real project makes you do. Deciding what to "
    "build (and often the biggest call: what the P2 shouldn't do at all — pair it with a "
    "Pi for the web stack rather than porting one). Learning the hardware — chasing "
    "datasheets, adding level shifters for 5V parts, and leaning on the logic analyzer as the "
    "court of final appeal. Building the capability — designing an interface someone can "
    "actually think in, translating reference drivers into Spin2, chasing performance, "
    "characterizing the real behavior (those measurements BECOME your product's spec). And "
    "finishing: document, post, announce — skip any of the three and you wasted the work. "
    "The key handoff: by the end of Part I you have parts that talk, a pin map, and a feel for "
    "the rates and deadlines — but NOT ONE COG assigned. That's exactly the raw material "
    "Part II works on. Tonight we're in Part II.")

# ---------------------------------------------------------------- Slide 3
bullet_slide(
    "The Big Idea: Computing in Space, Not Just in Time",
    [
        "A conventional MCU computes in TIME — one instruction stream, go faster",
        "An FPGA computes in SPACE — function laid out as parallel hardware",
        "The P2 sits between them: a coarse-grained spatial fabric",
        ("8 deterministic cogs + 64 smart pins you assign function to", 1),
        "Decompose WELL → parallel pipelines; throughput set by data RATE",
        "Decompose BADLY → the same silicon collapses to one cog doing everything,",
        ("the other seven idle", 1),
        "Decomposition is how you keep the design on the spatial side of that line",
    ],
    "CHAPTER 5 (≈90s). Here's the idea that makes the rest worth the effort. There are two "
    "ways a chip can compute. A normal microcontroller computes in TIME: one core, one "
    "instruction stream, and it does more by going faster or time-slicing. An FPGA is the "
    "opposite pole — it computes in SPACE: function laid out as actual parallel hardware, no "
    "single instruction stream. The P2 lives between them, closer to the spatial side than you'd "
    "expect — eight independent deterministic cogs and sixty-four smart pins form a "
    "coarse-grained spatial fabric. Here's the payoff and the warning in one: decomposed well, a "
    "P2 design behaves like spatial hardware — parallel pipelines whose throughput is set by "
    "the RATE data flows, not by how many instructions any one stage runs. Decomposed badly, the "
    "very same silicon collapses back into a slow sequential machine — one cog doing "
    "everything in turn while the other seven idle. That sentence is why this whole part exists. "
    "For P1 folks: you've been thinking this way all along — dedicate a cog, let it run. What "
    "the P2 changes is how MUCH fabric you have to spread onto.")

# ---------------------------------------------------------- FIGURE: space vs time
image_slide(
    "Computing in Time vs. in Space",
    "img/fig_spacetime.png",
    "FIGURE (show while covering the previous slide). The spectrum: a single-core MCU at the "
    "TIME pole (one instruction stream), an FPGA at the SPACE pole (function as parallel "
    "hardware), and the Propeller 2 sitting between them as a coarse-grained spatial fabric — "
    "8 cogs + 64 smart pins you assign function to. Leave this up while you deliver the "
    "decompose-well / decompose-badly point. You can merge this into the previous slide in "
    "Keynote if you'd rather have one slide.")

# ---------------------------------------------------------------- Slide 4
bullet_slide(
    "Object Shape Is DERIVED, Not Chosen",
    [
        "On the P2, your object set is not a style picked from a menu — it's derived",
        ("Change the buses, deadlines, or data rates → the correct object set changes", 1),
        "Vocabulary vs. grammar:",
        ("Objects (driver, policy, buffer, coordinator…) are the NOUNS", 1),
        ("The forces are the GRAMMAR — which nouns, how many, where the seams fall", 1),
        "Two axes, co-designed:",
        ("Logical — cohesion & coupling (decades of solid theory)", 1),
        ("Physical — allocation onto a finite, heterogeneous lattice", 1),
        "A smart pin can DELETE large amounts of code by absorbing the work into hardware",
    ],
    "CHAPTER 6 (≈90s). The central move of the whole part, stated plainly: on the P2 the "
    "shape of your object set is NOT a matter of taste picked from a menu. It's derived by "
    "reconciling a small number of physical and architectural forces. Change the buses, the "
    "deadlines, or the data rates and the correct object set changes with them. Why care? Because "
    "if decompositions were chosen by taste, the best you could do is collect examples and imitate "
    "the nearest one. Because they're DERIVED, you can learn the forces and produce a sound design "
    "for an app you've never seen. Separate two things: the OBJECTS — top-level app, device "
    "driver, semantic driver, policy layer, buffer, coordinator — are your vocabulary, the "
    "nouns. The FORCES are the grammar: the rules that decide which nouns, how many, and where the "
    "boundaries fall. Tonight is about the grammar. And the P2 adds a second axis classical "
    "software design doesn't have: a physical one — allocation onto eight cogs, sixty-four "
    "smart pins, one CORDIC, sixteen locks, bounded hub bandwidth. That physical axis isn't only a "
    "constraint, it's a design TOOL: a cog is the strongest encapsulation the silicon offers, and a "
    "smart pin can delete large amounts of code. 'Where's the boundary' and 'what hardware runs "
    "it' are ONE decision.")

# ---------------------------------------------------------------- Slide 5
bullet_slide(
    "The Four Forces That Do the Cutting",
    [
        "Three PRIMARY forces cut horizontally — who owns what, how pieces relate:",
        ("Force 1 — Who owns this wire?   (correctness)", 1),
        ("Force 2 — What does each seam promise?   (contracts)", 1),
        ("Force 3 — Where do two cadences meet?   (rate adapters)", 1),
        "One EMERGENT force falls out vertically once the first three draw the structure:",
        ("Force 4 — How high does each piece sit?   (layering)", 1),
        "Lead with the QUESTION each force asks — that's what you carry away",
        "For each force, the WHY matters more than the what",
    ],
    "CHAPTER 7 INTRO (≈60s). Four forces do the work. Three are PRIMARY — they cut the "
    "object set horizontally, deciding who owns what and how the pieces relate. The fourth is "
    "EMERGENT: it falls out vertically once the first three have drawn the structure. I'll lead "
    "each with the QUESTION it asks, because that question — asked of your OWN application "
    "— is the technique you take home. And for each one, the WHY matters more than the what: "
    "an engineer who knows why a force exists on the P2 will generalize it to hardware we never "
    "imagined; one who memorizes a rule eventually meets the case the rule didn't cover and applies "
    "it wrongly. So we'll dwell on reasons. The examples — a robot dog, some I2C buses — "
    "are a force in motion, never a rule to transplant.")

# ---------------------------------------------------------------- Slide 6
bullet_slide(
    "Force 1 — Who Owns This Wire?",
    [
        "A CORRECTNESS question — the only force that makes you flatly wrong, so it's first",
        "For each serialized, stateful resource (an I²C bus, a one-wire chain):",
        ("exactly ONE owning cog. The boundary traces the WIRE, not your feature list", 1),
        "Why: P2 pin outputs are OR'd — there is no hardware referee",
        ("Two cogs on one bus = guaranteed corruption, not 'a race you might lose'", 1),
        "A lock coordinates shared DATA; it can't un-corrupt a half-issued I²C frame",
        "The failure it prevents: the flat device list — every chip gets a sibling driver",
        "One owner, several shapes: many callers / broker cog / replicated bus",
    ],
    "FORCE 1 (≈90s). The first force asks a correctness question, not a style one. Of the "
    "four it's the only one that can make your program flatly WRONG rather than just inelegant — "
    "so it goes first. The question: for each serialized, stateful hardware resource — an I2C "
    "bus, a one-wire LED chain, a smart pin mid-transaction — which single cog owns it? The "
    "answer the force insists on: exactly one. One owner per resource, and the object boundary "
    "traces the WIRE, not your feature list. Why? It's physical. P2 pin outputs are OR'd together "
    "— there's no hardware referee arbitrating who drives the pin. If two cogs both drive SDA "
    "and SCL, they don't take polite turns; the outputs combine and the transaction corrupts. This "
    "isn't a race you might lose — a stateful protocol with two uncoordinated drivers is "
    "GUARANTEED to break. The chip gives you sixteen locks for shared DATA, but a lock can't "
    "un-corrupt a half-issued frame. The clean fix is structural: make the resource unshareable by "
    "giving it one owning object in one cog. The failure this prevents is the FLAT DEVICE LIST — "
    "every chip gets a driver, every driver a sibling off main(). The moment two cogs touch one bus "
    "you get silent, intermittent corruption that presents as flaky hardware, three layers from its "
    "cause. Note: one owner is the RULE; a single shared code image is just its cheapest ENCODING "
    "when there's one bus. Two identical buses? Give each its own state; keep one writer per bus.")

# ---------------------------------------------------------------- Slide 7
bullet_slide(
    "Force 2 — What Does Each Seam Promise?",
    [
        "Where two cogs meet, the 'contract' names the dependency — and draws the boundary",
        ("blocking call • latest-wins mailbox • ring buffer • published telemetry • fan-out", 1),
        "Stance: read a sensor continuously & publish its latest value — never on demand",
        "Every seam is really THREE planes, ranked by cost of getting them wrong:",
        ("Data plane → waste bandwidth (visible, recoverable)", 1),
        ("Control plane → corrupt state (an intermittent race)", 1),
        ("Event plane → miss a deadline (silent until the field)", 1),
        "Discipline: write the payload first, bump the signalling counter LAST (atomic long)",
    ],
    "FORCE 2 (≈90s). Once Force 1 scatters work across cogs, they have to exchange data. "
    "Force 2 asks: at each seam where two cogs meet, what does the exchange PROMISE? Does the "
    "sender wait? Does the receiver see the freshest value or every value? That promise is the "
    "CONTRACT, and choosing it IS a decomposition decision — the coupling you can tolerate "
    "decides where the boundary goes. Menu: a blocking call (tight coupling, caller waits on "
    "worst-case latency); a latest-wins mailbox (producer never waits, consumer always reads "
    "newest); a ring buffer (decouples rates, keeps every sample); published telemetry (one "
    "writer, many lockless readers); fan-out publication (one producer, many consumers, bulk "
    "frames). A design STANCE hides in the mailbox: for a sensor, read it continuously and publish "
    "its latest value — never reach out and read it at the moment you need it, because that "
    "couples your fast loop to the sensor's conversion latency. Now the sharpener: every seam is "
    "really THREE relationships superimposed, and I've ranked them by the cost of getting them "
    "wrong. Data plane — bulk movement; get it wrong, waste bandwidth, visible and "
    "recoverable. Control plane — commands and state; get it wrong, corrupt state, an "
    "intermittent race. Event plane — signalling and urgency; get it wrong, miss a deadline, "
    "and THAT one stays silent until the field. You spend design care in inverse order. The classic "
    "P2 abuse is building all three planes on one mechanism — polling a hub flag to deliver an "
    "urgent event. One discipline to memorize: publish-last — write the payload, bump the "
    "counter last; a single-long write is atomic, so a reader watching the counter can never catch "
    "a torn value. Costs nothing, removes a whole class of glitch.")

# ---------------------------------------------------------------- Slide 8
bullet_slide(
    "Force 3 — Where Do Two Cadences Meet?",
    [
        "The one a beginner most often misses — it corresponds to no chip, no parts-line item",
        "Devices live in different time domains: WS2812 (ns), servos (50 Hz), battery (1 Hz)…",
        "Whenever data crosses a cadence boundary, something must adapt the rate → an object",
        ("A software clock-domain crossing — glitches if you don't handle it on purpose", 1),
        "Two adapters fall out:",
        ("Sampler / buffer — fast producer, slow consumer (every sample? or freshest?)", 1),
        ("Slew / easing engine — a discrete 'walk' step → a smooth servo trajectory", 1),
        "Collides with Force 1: one bus, many cadences → cooperative tasks in the owner cog",
    ],
    "FORCE 3 (≈90s). This is the force a beginner's instinct most often misses, because it "
    "corresponds to nothing you can point at — no chip, no line in a parts list. Devices live "
    "in different TIME DOMAINS. An LED chain wants nanosecond bit timing; servos want a smooth 50 "
    "Hz; a voice recognizer is polled lazily and stretches the clock; a battery reading matters "
    "about once a second; an ultrasonic echo happens when it happens. Force 3 asks: where does data "
    "cross from one cadence to another, and what has to sit at that crossing to reconcile the rates? "
    "Because whenever data crosses a cadence boundary, SOMETHING must adapt the rate — and that "
    "adapter is a distinct responsibility, so it's a distinct object. The P2 encourages putting "
    "different time domains on different cogs and pins — that's what they're for — but the "
    "instant you do, you've built the software equivalent of a clock-domain crossing, and like its "
    "hardware namesake it glitches if you don't handle it on purpose. Two adapters fall out. A "
    "SAMPLER or buffer, where a fast producer meets a slow consumer — and the whole design is "
    "one question about the consumer: every sample (a buffer) or only the freshest (a latest-wins "
    "slot)? And a SLEW or easing engine, where a discrete intent becomes a continuous stream — "
    "'walk' arrives once, but a servo can't take a step; it needs a smooth ramp at its own frame "
    "rate. Pull the ramp out of both the policy and the driver and both stay clean. Best part: this "
    "COLLIDES with Force 1. Several devices share one bus but want different cadences — Force 1 "
    "forbids splitting the bus across cogs, Force 3 says separate the cadences. Neither wins by "
    "cutting. The answer is cooperative tasks WITHIN the one owning cog, each at its own cadence, "
    "yielding at transaction boundaries. You only find that by holding two forces in tension.")

# ---------------------------------------------------------------- Slide 9
bullet_slide(
    "Force 4 — How High Does Each Piece Sit?",
    [
        "The emergent, VERTICAL consequence — answers 'how much code goes in one object?'",
        "Not a line count. Split where the UNIT changes, or the axis of change changes",
        "The canonical stack: bits on a wire → registers → physical units → behavior",
        ("each tier = one unit conversion, changes for exactly one reason", 1),
        "Underneath: Parnas's information hiding — decompose around what changes independently",
        "Hard P2 limit: cog RAM is tiny (512 longs; 496 usable) — layering isn't free",
        ("Fold adjacent tiers when memory is tight — but SAY SO; never fold across reasons", 1),
        "Prevents the 'driver' that mixes register pokes with behavior logic",
    ],
    "FORCE 4 (≈90s). The first three forces are horizontal — who owns what, how they "
    "talk. The fourth is the VERTICAL consequence that falls out once they've drawn the structure, "
    "which is why it's emergent, not primary. It answers the question every programmer eventually "
    "asks: how much code goes in one object? The honest answer is NOT a line count and not a "
    "component count. It's: split where the unit changes, or where the axis of change changes. "
    "Stack the objects so each tier does exactly one unit conversion and changes for exactly one "
    "reason. The canonical stack climbs from bits on a wire, to device registers, to physical units "
    "— millimeters, degrees, millivolts — to behavior. Each tier speaks a different unit "
    "than the one below, and that change of unit IS the seam. The principle underneath is old and "
    "durable: Parnas's information hiding — decompose around the things that change "
    "independently, not around processing steps. Two pieces that always change together for the "
    "same reason belong in one object; two that change for different reasons — a new chip vs a "
    "new behavior — belong apart, even in the same call chain. On the P2 this negotiates "
    "against a hard limit: cog-local memory is TINY — 512 longs, 496 usable. Unlimited layering "
    "isn't free; each tier costs a call and a little state. So default to one tier per unit "
    "conversion, with an explicit escape: when a cog is genuinely tight, fold two adjacent tiers "
    "— but SAY SO, and never fold two tiers that change for different reasons just to save "
    "space, because that quietly rebuilds the monolith. The failure it prevents: a 'driver' that "
    "mixes register pokes with behavior, so swapping the IMU chip forces you to re-test the walk "
    "cycle. And a word on reconciling: the real skill isn't applying each force, it's reconciling "
    "them — they pull against each other. Hold them together, let them argue, and let the "
    "hardware and the hardest deadline win.")

# ---------------------------------------------------------------- Slide 10
bullet_slide(
    "Completing & Judging the Cut",
    [
        "The forces build a tree; some objects live ACROSS it — name them or they smear:",
        ("safety override • external-interface translator • config store • test seams • lifecycle sequencer", 1),
        "Keep a RESOURCE BUDGET as you derive — a blank row is a resource you forgot",
        ("'Running out of cogs' = the design is too COUPLED. Re-cut, don't cram", 1),
        ("Every cog needs a one-sentence reason: determinism / ownership / blocking I/O / throughput", 1),
        "Judge the cut — four tools of increasing sharpness:",
        ("coupling (countable longs) → change-coupling → back-pressure (min-cut) → observability", 1),
        "A decomposition is REVISABLE — expect to dial it in against real silicon",
    ],
    "CHAPTER 8 (≈100s). The four forces build a clean tree — but a real app needs objects "
    "that don't live IN the tree, they live ACROSS it. Five recur: a safety override (a privileged "
    "supervisor above policy — low-battery cutoff, e-stop); an external-interface translator "
    "(quarantine a vendor's vocabulary behind one seam); a configuration store (per-unit trim and "
    "pin maps in DATA, so identical firmware runs on every board); testability seams (each object "
    "exercised standalone on real hardware — the seam you can test at is the seam you should "
    "cut at); and a lifecycle sequencer (power before buses, chip awake before you actuate it). "
    "These have to be EXPLICIT on the P2 because cogs are independent — a hung cog won't stop "
    "driving its pins, and init order isn't implied by your call structure. Place them AFTER the "
    "tree is drawn. Next, keep a RESOURCE BUDGET as you derive — an allocation table, not a "
    "report you write afterward. Eight cogs, sixty-four pins, sixteen locks, one CORDIC, bounded "
    "hub bandwidth, 512 longs per cog. The budget earns its keep with one sharp signal: 'running "
    "out of cogs' is the P2's concrete way of telling you the design is too COUPLED — so re-cut, "
    "don't cram. And a check before you run out: every cog you assign needs a one-sentence reason "
    "from a short list — determinism, resource ownership, blocking I/O, or throughput. Can't say "
    "it? Fold it in. Finally, JUDGING — turning 'that seems cleaner' into something checkable. "
    "Four tools, increasing sharpness: coupling as a countable integer (count the longs that cross "
    "a boundary); change-coupling, the sharpest — what must CHANGE together (the dangerous case "
    "is dynamic change-coupling crossing a cog boundary, which the hardware expresses as jitter and "
    "races); back-pressure as a min-cut (draw the boundary where the least, weakest coupling crosses "
    "the cheapest channel); and observability — can you WATCH it run? On the P2 that's nearly "
    "free: aim a separate cog at a lock-free seam and observe WITHOUT perturbing. Last honesty: a "
    "decomposition isn't right just because you balanced the forces carefully once. It's a "
    "hypothesis; building the thing tests it. Expect to dial it in — that's the method working.")

# ---------------------------------------------------------------- Slide 11
bullet_slide(
    "The Method in Action — First-Contact Procedure",
    [
        "Deliberately inverts top-down: start from the hardware edge & the timing, not the data model",
        ("1. Enumerate the wires  (always)", 1),
        ("2. Triage against smart pins — a pin can DELETE large amounts of code  (skip if no pin mode fits)", 1),
        ("3. Assign owners — one cog, one transport per bus group  (always)", 1),
        ("4. List the cadences   5. Resolve same-bus rate conflicts", 1),
        ("6. Draw the seams (per plane)   7. Layer each branch (one tier per unit)", 1),
        ("8. Place the cross-cutting objects   9. Reconcile against budget & deadline  (always)", 1),
        "The procedure is FRACTAL — run it again inside a cog that owns a bus",
    ],
    "CHAPTER 9 PROCEDURE (≈60s). We have the forces, the cross-cutting objects, the budget, "
    "and the judging tools. The last thing you need is the ORDER to apply them — because the "
    "forces are orthogonal but the work isn't; you can't pick a seam's contract before you know "
    "where the cog boundaries are. The procedure deliberately INVERTS the classic top-down "
    "approach: you don't start from the data model, you start from the hardware edge and the timing "
    "budget and let the structure fall out. Nine steps — the spine steps always run, the "
    "others tell you when you can skip. Enumerate the wires. Triage against the smart pins — "
    "this is the physical axis used as a tool, a pin DELETES large amounts of code. Assign owners, one "
    "cog and one transport per bus group. List the cadences; resolve same-bus rate conflicts with "
    "cooperative tasks. Draw the seams per plane. Layer each branch, one tier per unit conversion. "
    "Place the cross-cutting objects — and name the ones you don't need. Reconcile against the "
    "budget and the hardest deadline. One more property: it's FRACTAL — after the top-level "
    "pass, run the very same routine INSIDE a cog that owns a bus; it has its own cadences, seams, "
    "and layers. When you're done you hold two things: the object-and-cog set, and the budget that "
    "proves it fits. Judge it with the four tools before you write a line.")

# ---------------------------------------------------------------- Slide 12
bullet_slide(
    "Watch It Run: A Walking Robot",
    [
        "Input is ONLY the hardware — nothing about the objects is given:",
        ("I²C bus 1: 13 servos + IMU + battery ADC, behind a ~50 Hz motion deadline", 1),
        ("I²C bus 2: one voice module that clock-stretches, polled slowly", 1),
        ("3 discretes: WS2812 chain, buzzer, ultrasonic ping/echo — smart pins carry the timing", 1),
        "The object set FALLS OUT: 3 cogs (orchestrator / body-control / I/O), ~8 smart pins,",
        ("0 locks (single-writer atomic telemetry). It fits, with cogs to spare — a min-cut", 1),
        "Same protocol (I²C) lands on two cogs with two transport shapes — topology decided it",
        "CARRY THE METHOD, NEVER THE MAP — a second app (streaming) gives a different answer",
    ],
    "CHAPTER 9 EXAMPLE (≈100s). Let's watch the whole method run once, end to end, on a small "
    "walking robot — a quadruped dog. The one thing that matters most: this is ONE "
    "application's answer, shown to make the method visible — it is NOT a template. Read for "
    "the moves, never the result. The only input is the hardware: I2C bus 1 carries thirteen servos "
    "through a PWM chip, plus an IMU and a battery ADC, all behind a hard ~50 Hz motion deadline. "
    "I2C bus 2 carries one voice module that clock-stretches and is polled slowly. And three "
    "discrete signals — an addressable LED chain, a buzzer, an ultrasonic ping-and-echo. "
    "Nothing about the objects is given; we derive them. Triage: the three discretes each map onto "
    "a smart-pin mode that carries the timing, so no cog bit-bangs them — and because the pins "
    "carry the jitter, all three collapse onto ONE non-blocking I/O cog. The two I2C buses are "
    "multi-byte stateful protocols — no smart pin can own I2C — so they need software "
    "owners. Bus 1's three devices sit behind one body-control cog with a single shared transport. "
    "Bus 1 serves three cadences — servos 50, IMU 100, battery 1 Hz — so three cooperative "
    "tasks inside that one cog. Seams: a latest-wins command mailbox (publish-last), lock-free "
    "published telemetry, event-plane freshness counters — nothing blocks. The motion branch "
    "layers into four tiers by unit conversion: PWM register driver, servo pulse/channel, leg "
    "inverse-kinematics, gait policy. Cross-cutting: a critical-battery hard-halt above policy, a "
    "voice-vocabulary translator at the edge, a per-joint trim store, per-layer bring-up tests, and "
    "the orchestrator owning launch order. Reconcile: three cogs of eight, about eight smart pins, "
    "ZERO locks — it fits with cogs to spare, and it's a min-cut. Notice what happened: we "
    "never started from a parts list and grabbed a template. We started from the wires and the "
    "timing, ran the forces in order, and the object set fell out — including the same protocol "
    "(I2C) on two cogs with two different state models, decided purely by sharing topology. The "
    "book runs a SECOND app — a fast image sensor streaming to two displays — through the "
    "identical nine steps and gets a totally different answer: a genuine FIFO pipeline, a decimator "
    "that sets the rate, and hub BANDWIDTH as the binding budget instead of cogs. Same method, "
    "different map. That's the whole claim in one line: carry the method, never the map.")

# ------------------------------------------------------ FIGURE: robot object map
image_slide(
    "The Walking Robot: Object-and-Cog Map",
    "img/fig_robot.png",
    "FIGURE (the visual payoff of the robot derivation). This is the object set the nine steps "
    "PRODUCED — not a template. Walk it left to right: the orchestrator/sequencer up top owns "
    "launch order and the safety supervisor; COG A (body-control) owns I2C bus 1 with its "
    "four-tier stack — gait policy, leg IK, servo semantics, PWM register driver — plus the "
    "per-unit trim store hanging off it; COG B (I/O cog) non-blocking-multiplexes the discretes "
    "and the voice bus. The labeled arrows are the seams by plane — CONTROL, DATA, EVENT. Point "
    "out: three cogs of eight, zero locks, and the SAME protocol (I2C) on two cogs with two "
    "transport shapes. Read it for the moves, not the result.")

# ------------------------------------------------------ FIGURE: streaming pipeline
image_slide(
    "A Second Application, a Different Answer",
    "img/fig_streaming.png",
    "FIGURE (the strongest evidence: same method, different map). Run the SAME nine steps on a "
    "fast image sensor streaming to two displays and you get something the robot never showed: a "
    "genuine FIFO pipeline. The tile sensor pours ~1,300 fps into a capture FIFO; a single "
    "DECIMATOR sets the rate — plain decimation or delta-compositing, lifting 8x8 toward 32x32 "
    "— then fans out through a FIFO per display to OLED and HDMI at ~60 fps. Three things "
    "differ from the robot: the producer gets its own cog for DETERMINISM (not bus ownership); "
    "the rate adapter is a pipeline, not in-cog cooperative tasks; and the binding budget is HUB "
    "BANDWIDTH, not cogs. None of the robot's boundaries carried over — the PROCEDURE did. "
    "That is the whole claim in one line: carry the method, never the map.")

# ---------------------------------------------------------------- Slide 13
bullet_slide(
    "Questions",
    [
        "Ask me anything on:",
        ("Computing in space vs. time — and where the P2 really sits", 1),
        ("Any of the four forces — ownership, seams, cadences, layering", 1),
        ("The resource budget and the four judging tools", 1),
        ("The walking-robot derivation — or the streaming-pipeline contrast", 1),
        "The full treatment (plus appendices & glossary) is in the released Guide",
    ],
    "QUESTIONS (open-ended). Take questions here. Likely ones and short answers: "
    "— 'Isn't three cogs wasteful when I have eight?' A healthy design ships with a cog or two "
    "in reserve; each cog in use should carry its one forcing sentence. Reserve beats reflexively "
    "filling all eight. "
    "— 'When do I use a lock vs. atomic publish?' Locks coordinate shared DATA; for a "
    "single-writer/many-reader telemetry hand-off you often need NO lock — publish-last on an "
    "atomic long. "
    "— 'What if a smart pin COULD do it but bit-banging is faster?' That's a signature to "
    "check, not a verdict — override it if you must, but record why and what would reverse the "
    "decision (the streaming example's OLED is exactly this). "
    "— 'How do I know my cut is good?' Run the four tools — and remember a decomposition "
    "is revisable; you dial it in against real silicon. "
    "If a question goes deep, point them to the released Guide — the appendices cover the "
    "FPGA borrowing and the decomposition literature.")

# ---------------------------------------------------------------- Slide 14
bullet_slide(
    "Let's Discuss — Your Decomposition Experiences",
    [
        "I'd love to hear how YOU decompose P2 designs — let's compare notes:",
        ("How do YOU decide what goes on which cog? What's your first move?", 1),
        ("A decomposition that FOUGHT you — what did you have to re-cut, and why?", 1),
        ("Does the 'derived, not chosen' framing match how it actually feels?", 1),
        ("Anything that guides YOUR decomposition we haven't named — a force we missed?", 1),
        "And if an example really fits the book's approach, it might inform a future edition — credited",
    ],
    "DISCUSSION (the real payoff — keep it going as long as it's live). This is why I wanted to "
    "present tonight rather than just post the book. The Guide's method is derived from a dozen "
    "real projects, but many of YOU have decomposed P2 designs I've never seen — and I'd love "
    "to learn from that experience tonight. Prompts to get us going: How do you decide "
    "what goes on which cog — what's your very first move on a new hardware mix? Where have you "
    "run out of cogs, and looking back, was it genuinely at capacity or was it too-coupled in "
    "disguise? Tell me about a decomposition that FOUGHT you — the seam that looked clean on "
    "paper and turned awkward in code, the cut you had to redraw. Any case where a smart pin let "
    "you delete large amounts of code you'd otherwise have spent a cog on. And where a blocking call between "
    "cogs quietly serialized something that was meant to run in parallel. Big-picture: does "
    "'object shape is derived, not chosen' match how it actually feels when you design, or does it "
    "feel more like taste to you? And the one I care about most: is there anything that guides your "
    "decomposition we haven't named tonight — a force we missed? Be honest — pushback makes the "
    "next edition better. I'll "
    "capture the strong examples; and if any really fit the book's approach, then with your "
    "permission they might inform a future edition, credited to you. Thank you — who wants to start?")

name = ("thinking-in-p2-decomposition-noimg-test.pptx" if NO_IMAGES
        else "thinking-in-p2-decomposition.pptx")
for i, slide in enumerate(prs.slides, start=1):
    if i > 1:                       # title slide (1) carries its own credit block
        add_footer(slide, i)

for slide in prs.slides:
    apply_gradient_bg(slide)

out = __file__.rsplit("/", 1)[0] + "/" + name
prs.save(out)
wire_notes_master(out)   # Keynote: inject <p:notesMasterIdLst> python-pptx omits
print("wrote", out, "with", len(prs.slides._sldIdLst), "slides",
      "(NO_IMAGES)" if NO_IMAGES else "")
